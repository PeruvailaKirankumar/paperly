import os
from typing import List
from PyPDF2 import PdfReader
from pptx import Presentation
import io
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path


class DocumentProcessor:
    """Handles extraction of text from various document formats."""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF file."""
        try:
            reader = PdfReader(file_path)
            text = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
            return "\n\n".join(text)
        except Exception as e:
            raise ValueError(f"Error reading PDF file: {str(e)}")
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_text.append(f"--- Slide {slide_num} ---")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the header
                    text.append("\n".join(slide_text))
            
            return "\n\n".join(text)
        except Exception as e:
            raise ValueError(f"Error reading PPTX file: {str(e)}")
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            with zipfile.ZipFile(file_path) as zip_file:
                # Extract text from document.xml
                if 'word/document.xml' in zip_file.namelist():
                    xml_content = zip_file.read('word/document.xml')
                    root = ET.fromstring(xml_content)

                    # Extract all text from paragraph elements
                    text_content = []
                    for elem in root.iter():
                        if elem.tag.endswith('}p'):  # Paragraph
                            paragraph_text = []
                            for child in elem:
                                if child.tag.endswith('}t'):  # Text
                                    if child.text:
                                        paragraph_text.append(child.text)
                            if paragraph_text:
                                text_content.append(' '.join(paragraph_text))

                    return '\n\n'.join(text_content)
                else:
                    raise ValueError("Invalid DOCX file: no document.xml found")
        except Exception as e:
            raise ValueError(f"Error reading DOCX file: {str(e)}")

    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Extract text from TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                raise ValueError(f"Error reading TXT file: {str(e)}")
        except Exception as e:
            raise ValueError(f"Error reading TXT file: {str(e)}")

    @staticmethod
    def extract_text_from_ppt(file_path: str) -> str:
        """Extract text from PPT file (older format).
        Note: For better PPT support, consider using python-pptx or comtypes."""
        # PPT files (older format) are more complex to parse
        # For now, we'll return an error message suggesting conversion
        raise ValueError(
            "Legacy PPT format not fully supported. "
            "Please convert to PPTX format for best results. "
            "Alternatively, install additional libraries like comtypes (Windows only)."
        )
    
    @staticmethod
    def extract_text(file_path: str) -> str:
        """Extract text from file based on extension."""
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        if ext == ".pdf":
            return DocumentProcessor.extract_text_from_pdf(file_path)
        elif ext == ".pptx":
            return DocumentProcessor.extract_text_from_pptx(file_path)
        elif ext == ".ppt":
            return DocumentProcessor.extract_text_from_ppt(file_path)
        elif ext == ".docx":
            return DocumentProcessor.extract_text_from_docx(file_path)
        elif ext == ".txt":
            return DocumentProcessor.extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    @staticmethod
    def get_supported_extensions() -> List[str]:
        """Return list of supported file extensions."""
        return [".pdf", ".pptx", ".docx", ".txt"]
